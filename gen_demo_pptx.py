#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BLUE = RGBColor(0x1a, 0x3a, 0x5c)
LIGHT_BLUE = RGBColor(0x2c, 0x5f, 0x8d)
ACCENT_ORANGE = RGBColor(0xe8, 0x6a, 0x3a)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x2e, 0x7d, 0x32)
LIGHT_GREEN = RGBColor(0xe8, 0xf5, 0xe9)
LIGHT_BLUE_BG = RGBColor(0xe3, 0xf2, 0xfd)
RED_LIGHT = RGBColor(0xff, 0xf0, 0xf0)
RED_BORDER = RGBColor(0xef, 0x9a, 0x9a)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, text, left, top, width, height,
       size=16, bold=False, color=DARK_GRAY,
       align=PP_ALIGN.LEFT, name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    return box

def header_bar(slide, title):
    s = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(29.7), Cm(2.5))
    s.fill.solid()
    s.fill.fore_color.rgb = DARK_BLUE
    s.line.fill.background()
    tb(slide, title, Cm(1), Cm(0.5), Cm(27), Cm(1.8), 22, True, WHITE)

def scene_label(slide, num, title):
    box = slide.shapes.add_shape(1, Cm(1), Cm(3), Cm(3.5), Cm(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_ORANGE
    box.line.fill.background()
    tb(slide, "Scene " + str(num), Cm(1.2), Cm(3.1), Cm(3.3), Cm(0.8), 12, True, WHITE)
    tb(slide, title, Cm(5), Cm(3.1), Cm(23), Cm(0.8), 15, True, DARK_BLUE)

def narration(slide, text):
    box = slide.shapes.add_shape(1, Cm(1), Cm(4.2), Cm(27.7), Cm(6.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    tb(slide, "[Narration]", Cm(1.3), Cm(4.3), Cm(3), Cm(0.6), 11, True, LIGHT_BLUE)
    tb(slide, text, Cm(1.3), Cm(4.8), Cm(27), Cm(5.5), 13, False, DARK_GRAY)

def interaction(slide, text):
    box = slide.shapes.add_shape(1, Cm(1), Cm(11.2), Cm(27.7), Cm(3))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE_BG
    box.line.color.rgb = LIGHT_BLUE
    tb(slide, "[Interaction]", Cm(1.3), Cm(11.3), Cm(4), Cm(0.6), 11, True, LIGHT_BLUE)
    tb(slide, text, Cm(1.3), Cm(11.7), Cm(27), Cm(2.3), 13, False, DARK_BLUE)

def page_num(slide, cur, total):
    tb(slide, str(cur) + " / " + str(total), Cm(25), Cm(19.5), Cm(4), Cm(1),
       11, False, MID_GRAY, PP_ALIGN.RIGHT)

TOTAL = 12

prs = Presentation()
prs.slide_width = Cm(29.7)
prs.slide_height = Cm(21)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BLUE)
tb(slide, "AI商务社交", Cm(2), Cm(5), Cm(25), Cm(4),
   48, True, WHITE, PP_ALIGN.CENTER)
tb(slide, "Demo Script", Cm(2), Cm(9.5), Cm(25), Cm(2),
   28, False, RGBColor(0xcc, 0xcc, 0xcc), PP_ALIGN.CENTER)
tb(slide, "Sales/Pre-sales Demo | Duration: 15-20 min",
   Cm(2), Cm(14), Cm(25), Cm(1.5), 12, False, RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)
page_num(slide, 1, TOTAL)

# Slide 2: How to use
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "How to Use This Script")
tb(slide, "[Presenter Role]", Cm(1.5), Cm(3.2), Cm(26), Cm(0.6), 14, True, DARK_BLUE)
tb(slide, "You are the facilitator. The client CEO is the protagonist. "
   "The client will role-play as an EMBA alumni experiencing how AI helps discover networking opportunities.",
   Cm(1.5), Cm(3.7), Cm(26), Cm(2), 13, False, DARK_GRAY)
tb(slide, "[How to Use]", Cm(1.5), Cm(5.8), Cm(26), Cm(0.6), 14, True, DARK_BLUE)
lines = [
    "1. Prepare in advance: have PPT/mockup screenshots ready for each scene",
    "2. Follow the sequence: 5 scenes, 2-4 min each, total 15-20 min",
    "3. Read [Narration] aloud: this is what you say",
    "4. Watch [Interaction] points: pause here for client engagement",
    "5. Show product interface when indicated"
]
top = Cm(6.4)
for line in lines:
    tb(slide, line, Cm(1.5), top, Cm(26), Cm(0.9), 13, False, DARK_GRAY)
    top += Cm(1.0)
page_num(slide, 2, TOTAL)

# Slide 3: Scene Setup
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Scene Setup (Opening)")
scene_label(slide, 0, "Opening Introduction")
narration(slide,
    "Mr. Wang, welcome to experience our product. Today I'll take you through a complete demo.\n\n"
    "Your role: A 02nd Cohort EMBA alumnus, currently working on an AI-assisted diagnosis "
    "startup project, looking to collaborate with medical resources.\n\n"
    "Let's start by looking at what your alumni homepage looks like in the product.")
interaction(slide, "At opening, confirm client's basic profile: industry, background, "
    "collaboration intent. Adjust the role-play scenario accordingly.")
page_num(slide, 3, TOTAL)

# Slide 4: Scene 1 - Profile
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Scene 1: AI Analysis of Alumni Profile")
scene_label(slide, 1, "Alumni Profile -- Meet Chen Daming")
narration(slide,
    "First, let's see: if you wanted to learn more about an alumnus -- say Chen Daming, "
    "partner at a medical investment fund -- how would you do it traditionally?\n\n"
    "In our product, you just need to look at the AI analysis.")
interaction(slide, "After showing the product interface (Alumni AI Analysis), ask the client:\n"
    '"In traditional products, how do you usually learn about an alumnus? How long does it take?"')
page_num(slide, 4, TOTAL)

# Slide 5: Scene 1 - Value
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Scene 1: Value of Alumni Profile")

box1 = slide.shapes.add_shape(1, Cm(1), Cm(4.2), Cm(13), Cm(8))
box1.fill.solid()
box1.fill.fore_color.rgb = RED_LIGHT
box1.line.color.rgb = RED_BORDER
tb(slide, "Traditional Product:", Cm(1.3), Cm(4.4), Cm(12), Cm(0.6), 12, True, RGBColor(0xc6, 0x28, 0x28))
tb(slide, "Name: Chen Daming\nClass: 02nd EMBA\nCompany: Medical Fund\nTitle: Partner",
   Cm(1.3), Cm(5.1), Cm(12), Cm(6.5), 14, False, DARK_GRAY)

box2 = slide.shapes.add_shape(1, Cm(15.5), Cm(4.2), Cm(13.2), Cm(8))
box2.fill.solid()
box2.fill.fore_color.rgb = LIGHT_GREEN
box2.line.color.rgb = GREEN
tb(slide, "Our Product + AI:", Cm(15.8), Cm(4.4), Cm(12), Cm(0.6), 12, True, GREEN)
tb(slide, "Tags: Investment, Industry Resources, Early-stage Projects\n"
   "Supply & Demand: Has=Investment resources, Needs=AI early-stage projects\n"
   "Match Analysis: 4 connection points listed\n"
   "Ice-breaker: Talk data first, don't mention funding",
   Cm(15.8), Cm(5.1), Cm(12), Cm(6.5), 14, False, DARK_GRAY)

tb(slide, '[Narration] This is what an alumni directory should look like -- '
   'not an address book, but an intelligent relationship network.',
   Cm(1), Cm(13.2), Cm(27.7), Cm(1.5), 13, False, DARK_BLUE)
page_num(slide, 5, TOTAL)

# Slide 6: Scene 2 - Traditional
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Scene 2: Two Ways to Find People")
scene_label(slide, "2a", "Traditional Search")

box = slide.shapes.add_shape(1, Cm(1), Cm(4.2), Cm(27.7), Cm(7))
box.fill.solid()
box.fill.fore_color.rgb = RED_LIGHT
box.line.color.rgb = RED_BORDER
tb(slide, "[Traditional Search -- Mr. Wang's Path]", Cm(1.3), Cm(4.4), Cm(26), Cm(0.6),
   13, True, RGBColor(0xc6, 0x28, 0x28))
tb(slide,
   "1. Search 'medical' in alumni directory -> 50 names\n"
   "2. Check each profile one by one to judge relevance\n"
   "3. Pick 3 that look promising, send: 'Hello, I am Mr. Wang, may I ask about collaboration opportunities'\n"
   "4. Wait 3 days -> 1 reply: 'Sorry, we mainly look at late-stage projects, you're still early'\n\n"
   "-> A whole day wasted.",
   Cm(1.3), Cm(5.1), Cm(26), Cm(5.5), 14, False, DARK_GRAY)
page_num(slide, 6, TOTAL)

# Slide 7: Scene 2 - AI
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Scene 2: Two Ways to Find People")
scene_label(slide, "2b", "AI Networking Assistant")

box = slide.shapes.add_shape(1, Cm(1), Cm(4.2), Cm(27.7), Cm(8))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_GREEN
box.line.color.rgb = GREEN
tb(slide, "[AI Assistant -- Mr. Wang's Path]", Cm(1.3), Cm(4.4), Cm(26), Cm(0.6),
   13, True, GREEN)
tb(slide,
   "1. Open AI Assistant, type: 'Working on an AI-assisted diagnosis project, looking for alumni with medical resources'\n"
   "2. AI finds 3 high-match alumni, each with match score and reason:\n"
   "   - Chen Daming (02nd EMBA | Medical Fund) -- Match: 95%\n"
   "   - Liu Fang (15th MBA | Top Hospital) -- Match: 88%\n"
   "   - Zhao Qiang (08th EMBA | Medical Group) -- Match: 82%\n"
   "3. AI has drafted the first message -- just confirm and send",
   Cm(1.3), Cm(5.1), Cm(26), Cm(6.5), 13, False, DARK_GRAY)

tb(slide, '[Narration] Not 50 people to choose from yourself, but directly telling you why these 3 are worth knowing.',
   Cm(1), Cm(13), Cm(27.7), Cm(1.5), 13, False, DARK_BLUE)
interaction(slide,
    '"Mr. Wang, if you were the user, how would you describe your need for AI medical collaboration to the AI? '
    'Just say it naturally."')
page_num(slide, 7, TOTAL)

# Slide 8: Scene 3 - Proactive
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Scene 3: AI-Proactively Discovered Opportunity")
scene_label(slide, 3, "Supply-Demand Match Push -- Proactive Discovery")
narration(slide,
    "The previous scenes were you actively searching. But the more valuable part is -- "
    "AI proactively discovering opportunities you didn't know about.\n\n"
    "One day you open the product and see a push notification...")
interaction(slide,
    'After showing product interface (Supply-Demand Match Push), ask:\n\n'
    '"Mr. Wang, look at this -- Zhang Wei posted this need, and AI pushed it to you proactively. '
    'You might not have known Zhang Wei was looking for AI pathology collaboration, '
    "but AI knew, and AI knew you have this capability.\n"
    'If you saw this push, would you want to click in?"')
page_num(slide, 8, TOTAL)

# Slide 9: Scene 4 - Group chat
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Scene 4: AI Connecting People in Alumni Group Chat")
scene_label(slide, 4, "EMBA Alumni Group -- AI in the Group Chat")
narration(slide,
    "Now for a more interesting scenario -- AI in an alumni group chat.\n"
    "The EMBA group has 127 people. Everyone posts messages and discusses collaboration every day.\n"
    "Mr. Wang asks in the group: 'Any alumni working on medical device CDMO? "
    "We have a project looking for a reliable manufacturer, requires NMPA certification.'")
interaction(slide,
    'After showing product interface, ask:\n\n'
    '"Mr. Wang, if you were Li Na -- you just joined the group, '
    'and received this AI notification telling you someone is looking for CDMO collaboration. '
    'How would you feel?"')
tb(slide, '[Narration] Li Na can choose whether to respond. If she says yes, AI writes the reply for her.',
   Cm(1), Cm(13.5), Cm(27.7), Cm(1.5), 13, False, DARK_BLUE)
page_num(slide, 9, TOTAL)

# Slide 10: Scene 5 - Activity
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Scene 5: AI in Event Scenarios")
scene_label(slide, 5, "Event Detail Page -- Attendee Recommendations")
narration(slide,
    "The last scenario -- events. This is another important module in our product.\n"
    "Let's say there's an AI Medical Application Innovation Forum next week.\n\n"
    "Traditional event logic: you register, then you figure out who to talk to.\n"
    "Our product reverses this: before you register, AI tells you who you should know at this event "
    "and why.")
tb(slide, '[Narration] An event is no longer a one-time social occasion, '
   'but the starting point of an ongoing relationship.',
   Cm(1), Cm(13.5), Cm(27.7), Cm(1.5), 13, False, DARK_BLUE)
page_num(slide, 10, TOTAL)

# Slide 11: Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
header_bar(slide, "Summary: Core Product Logic")

# Summary table
headers = ["Module", "Traditional", "Our Product"]
rows = [
    ["Alumni Directory", "Store contacts", "AI finds who you should know"],
    ["Social Events", "Event ends, connection ends", "AI helps ongoing follow-up"],
    ["Supply-Demand", "Bulletin board", "AI matches latent supply-demand"],
    ["AI Assistant", "Must learn how to use", "Just say one sentence"],
]
col_widths = [Cm(5), Cm(10), Cm(13.7)]
col_starts = [Cm(1.5), Cm(7), Cm(17.5)]
row_h = Cm(1.4)
top = Cm(3.5)

for i, h in enumerate(headers):
    box = slide.shapes.add_shape(1, col_starts[i], top, col_widths[i], Cm(1))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BLUE
    box.line.fill.background()
    tb(slide, h, col_starts[i] + Cm(0.2), top + Cm(0.2), col_widths[i], Cm(1),
       14, True, WHITE)

top += row_h
for row in rows:
    for i, cell in enumerate(row):
        box = slide.shapes.add_shape(1, col_starts[i], top, col_widths[i], row_h)
        if i == 1:
            box.fill.solid()
            box.fill.fore_color.rgb = RED_LIGHT
            box.line.color.rgb = RED_BORDER
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_GREEN
            box.line.color.rgb = GREEN
        tb(slide, cell, col_starts[i] + Cm(0.2), top + Cm(0.2), col_widths[i], row_h,
           13, False, DARK_GRAY)
    top += row_h

tb(slide, "Alumni is the foundation, AI is the engine, events are the catalyst, supply-demand is the driver.",
   Cm(1.5), Cm(10.5), Cm(26), Cm(2),
   16, True, DARK_BLUE, PP_ALIGN.CENTER)
tb(slide, 'Not helping you find people, but helping you discover connections you didn\'t know you needed.',
   Cm(1.5), Cm(12.5), Cm(26), Cm(2),
   14, False, MID_GRAY, PP_ALIGN.CENTER)
page_num(slide, 11, TOTAL)

# Slide 12: Closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BLUE)
tb(slide, "Mr. Wang, what do you think about this direction?",
   Cm(2), Cm(5.5), Cm(25), Cm(3),
   26, True, WHITE, PP_ALIGN.CENTER)
tb(slide, "Is there any particular scenario that resonates with you?\nWe can expand further on any of these.",
   Cm(2), Cm(9.5), Cm(25), Cm(3),
   17, False, RGBColor(0xcc, 0xcc, 0xcc), PP_ALIGN.CENTER)
tb(slide, "(Wait for client feedback, explore further based on client's interest)",
   Cm(2), Cm(14), Cm(25), Cm(2),
   12, False, RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)
page_num(slide, 12, TOTAL)

out = "/Users/wangsheng/luojing/Demo演示脚本.pptx"
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
