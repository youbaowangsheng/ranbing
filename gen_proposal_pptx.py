#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BLUE = RGBColor(0x1a, 0x3a, 0x5c)
LIGHT_BLUE = RGBColor(0x2c, 0x5f, 0x8d)
ACCENT_ORANGE = RGBColor(0xe8, 0x6a, 0x3a)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, text, left, top, width, height,
       size=16, bold=False, color=DARK_GRAY,
       align=PP_ALIGN.LEFT, name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    return box

def header_bar(slide, title):
    s = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(29.7), Cm(2.5))
    s.fill.solid()
    s.fill.fore_color.rgb = DARK_BLUE
    s.line.fill.background()
    tb(slide, title, Cm(1), Cm(0.5), Cm(27), Cm(1.8), 22, True, WHITE)

def content_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    header_bar(slide, title)
    top = Cm(3.5)
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
            indent = level * 0.8
        else:
            text = item
            level = 0
            indent = 0
        prefix = "  " if level > 0 else ""
        bullet_char = "  " if level > 0 else "  "
        display = prefix + text
        size = 13 if level > 0 else 15
        b = level == 0
        tb(slide, display, Cm(1.5 + indent), top, Cm(26), Cm(1.1),
           size, b, DARK_GRAY)
        top += Cm(1.15)
    if note:
        tb(slide, note, Cm(1.5), Cm(19.5), Cm(26), Cm(1.5),
           11, False, MID_GRAY)
    return slide

def two_col_slide(prs, title, lt, lb, rt, rb):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    header_bar(slide, title)

    box1 = slide.shapes.add_shape(1, Cm(1), Cm(3.2), Cm(13), Cm(1))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_BLUE
    box1.line.fill.background()
    tb(slide, lt, Cm(1.3), Cm(3.35), Cm(12.5), Cm(0.8), 13, True, WHITE)

    box2 = slide.shapes.add_shape(1, Cm(15.5), Cm(3.2), Cm(13), Cm(1))
    box2.fill.solid()
    box2.fill.fore_color.rgb = ACCENT_ORANGE
    box2.line.fill.background()
    tb(slide, rt, Cm(15.8), Cm(3.35), Cm(12.5), Cm(0.8), 13, True, WHITE)

    top = Cm(4.8)
    for b in lb:
        tb(slide, "  " + b, Cm(1.3), top, Cm(12.5), Cm(0.9), 13, False, DARK_GRAY)
        top += Cm(1.0)
    top = Cm(4.8)
    for b in rb:
        tb(slide, "  " + b, Cm(15.8), top, Cm(12.5), Cm(0.9), 13, False, DARK_GRAY)
        top += Cm(1.0)
    return slide

def quote_slide(prs, quote, attr=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    tb(slide, '"', Cm(1), Cm(4), Cm(3), Cm(3), 120, True, LIGHT_BLUE)
    tb(slide, quote, Cm(2.5), Cm(6), Cm(24), Cm(8), 22, False, DARK_BLUE)
    if attr:
        tb(slide, "-- " + attr, Cm(2.5), Cm(15), Cm(24), Cm(1.5), 14, False, MID_GRAY, PP_ALIGN.RIGHT)
    return slide

prs = Presentation()
prs.slide_width = Cm(29.7)
prs.slide_height = Cm(21)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BLUE)
tb(slide, "AI时代商务社交解决方案", Cm(2), Cm(5.5), Cm(25), Cm(4),
   40, True, WHITE, PP_ALIGN.CENTER)
tb(slide, '从"电子通讯录"到"智能关系网络运营平台"', Cm(2), Cm(10), Cm(25), Cm(2),
   18, False, RGBColor(0xcc, 0xcc, 0xcc), PP_ALIGN.CENTER)
tb(slide, "产品方案 v2.0 | 2026年5月", Cm(2), Cm(14), Cm(25), Cm(1.5),
   12, False, RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)

# Slide 2: Background
content_slide(prs, "时代背景：为什么现在需要重新思考商务社交", [
    ('传统商务社交在PC时代诞生，移动时代做了"App化"，核心逻辑从未改变', 0),
    ('用户路径: 注册 -> 完善资料 -> 搜索"某总" -> 发消息"您好" -> 石沉大海', 0),
    ('或者: 注册 -> 参加活动 -> 换名片 -> 加微信 -> 再无联系', 0),
    ("三个根本性问题：", 0),
    ("关系静态化 -- 通讯录里的名字不会主动产生价值", 1),
    ("信息严重滞后 -- 名片上的公司是三年前的，需求早已变化", 1),
    ("匹配极度低效 -- 发10条消息可能只有1条回复", 1),
])

# Slide 3: Traditional vs AI
two_col_slide(prs, "AI带来的结构性机会",
    "传统模式",
    ["人找信息", "用户主动搜索", '关系靠"刷脸"维护', "通讯录是死的", "匹配靠猜"],
    "AI模式",
    ["信息找人", "AI主动推送", "AI持续激活关系", "人脉图谱是活的", "匹配靠算法"])

# Slide 4: Vision
quote_slide(prs,
    '不是帮你找到人，而是帮你发现你不知道该找谁的连接。',
    "产品愿景")

# Slide 5: Architecture
content_slide(prs, "产品架构：三层结构", [
    ("第一层：动态人脉图谱", 0),
    ("  节点: 用户Profile、供需信息、活动轨迹、标签体系", 1),
    ("  边: 互动历史、共同属性、匹配度评分（AI计算）", 1),
    ("第二层：AI社交引擎", 0),
    ("  意图理解 / 关系发现 / 机会匹配 / 内容生成 / 持续学习", 1),
    ("第三层：用户交互界面", 0),
    ("  主页=AI今日推送  |  私信=AI辅助写作", 1),
], note='核心变化: 用户不是"去搜索"，而是"来看AI今天发现了什么"')

# Slide 6: Four modules
content_slide(prs, "四大功能模块：社交飞轮", [
    ("校友网络 -- 信任底座，同学关系天然带信任背书", 0),
    ("社交活动 -- 关系催化剂，让关系快速建立和加深", 0),
    ("供需广场 -- 匹配引擎，驱动产生实际合作", 0),
    ("AI人脉助手 -- 统一入口，降低用户使用门槛", 0),
    ("任何模块的数据回流 -> 增强人脉图谱 -> AI能力提升 -> 更多匹配", 0),
], note="四个模块不是四个独立功能，而是一个互相增强的飞轮")

# Slide 7: Module 1 - Alumni
content_slide(prs, "模块一：校友名录（信任底座）", [
    ("多维度校友档案: 年级、班级、专业、EMBA/MBA、社团、创业经历", 0),
    ("自完善体系: 引导用户逐步完善资料，认证后获得信任标识", 0),
    ("校友社群: 按年级、专业、地域、兴趣组建社群，支持群内互动", 0),
    ("互动私信: 支持一对一私信，有已读回执，AI辅助写作", 0),
    ("AI自动标签化: 从文本里学习，自动打标签", 0),
    ('AI隐性关系发现: 发现"两个校友都在关注AI医疗"等隐性关联', 0),
    ("AI智能推荐: 基于人脉网络，推荐值得主动连接的校友", 0),
])

# Slide 8: Module 1 - Demo
content_slide(prs, "模块一：AI深度解读校友主页", [
    ("AI自动解析: 能力标签、供需分析、与你的匹配分析", 0),
    ("匹配分析示例:", 0),
    ("  同是02届EMBA，有共同圈子", 1),
    ("  他在找AI辅助诊断，你正好在这个方向", 1),
    ("  你的项目已有2家医院落地，符合他的偏好", 1),
    ("  他今天刚发动态关注这个赛道", 1),
    ("AI破冰建议: 先聊技术和临床验证数据，不要直接谈融资", 0),
], note="价值: 不是通讯录，而是一个智能关系网络")

# Slide 9: Module 2 - Activities
content_slide(prs, "模块二：社交活动（关系催化剂）", [
    ("活动发布与管理: 发起活动、设置报名门槛、签到管理", 0),
    ("活动日历: 按时间、行业、规模筛选活动", 0),
    ("活动推荐: 基于用户标签和关系网络推荐值得参加的活动", 0),
    ("同行识别: 报名同一个活动的校友里，谁可能是你想认识的", 0),
    ('活动后跟进: AI自动生成"你可能想跟进的人"列表', 0),
    ("AI在活动前/中/后的介入:", 0),
    ("  活动前: 预判哪些人值得提前搭话", 1),
    ("  活动中: 记录谁和谁交流过，识别深度连接可能性", 1),
    ("  活动后: 生成跟进列表，AI起草第一条消息", 1),
])

# Slide 10: Module 3 - Supply/Demand
content_slide(prs, "模块三：供需广场（匹配引擎）", [
    ('供需发布: 用户发布"我有什么"（供给）或"我在找什么"（需求）', 0),
    ("智能匹配: AI主动推送匹配的供需给最可能响应的人", 0),
    ('意向确认: AI先"预对话"确认意向，再转真人', 0),
    ("AI介入四步:", 0),
    ("  第一步: AI解析供需，把文本转化为结构化标签", 1),
    ("  第二步: AI发现匹配，A的供给恰好是B的需求", 1),
    ("  第三步: AI主动推送，不等用户来搜索", 1),
    ('  第四步: AI预对话确认意向，大幅减少"发出去没回复"', 1),
], note="价值最直接: 两个人因为供需匹配开始合作，这个价值是可验证的")

# Slide 11: Module 4 - AI Assistant
content_slide(prs, "模块四：AI人脉助手（统一入口）", [
    ("自然语言交互: 用户输入任意自然语言需求", 0),
    ("意图识别: 判断用户是想找资源、找人、找活动、还是问问题", 0),
    ("需求转化: 把模糊需求变成结构化的AI能处理的query", 0),
    ("智能分发: 根据意图，把需求导向最合适的功能模块", 0),
    ("用户只需要说一句话，不需要学习怎么用这个产品", 0),
    ("示例:", 0),
    ('  "我想找AI医疗的人聊聊" -> 导向校友推荐', 1),
    ('  "最近有什么活动？" -> 导向活动推荐', 1),
    ('  "帮我起草一段给李娜的消息" -> AI辅助写作', 1),
])

# Slide 12: AI Value
two_col_slide(prs, "AI价值如何显性化",
    "弱感知（无效做法）",
    ['搜索结果标注"AI排序"', '弹出通知"AI发现新机会"', '列表页标注"AI推荐"', "对话机器人回答问题"],
    "强感知（有效做法）",
    ['直接告诉你"这3个人为什么值得认识"', "明确展示匹配理由和关系链", "AI帮你起草第一条消息，只需确认", "AI在你还没意识到时就推送到面前"])

# Slide 13: Three principles
content_slide(prs, "强感知设计的三个原则", [
    ('原则1: AI每做一个动作，必须给用户一个"我为什么应该关心"的答案', 0),
    ("原则2: 关系链要可视化 -- 用户能看到这段关系的来龙去脉", 0),
    ("原则3: 下一步动作要极简 -- 用户不需要从零打字，只需要确认或修改", 0),
])

# Slide 14: Homepage
content_slide(prs, "产品主页: AI今天发现了什么", [
    ("[通知] 供需匹配 -- 张伟刚发了AI病理合作需求，你们是同班同学，匹配度94%", 0),
    ("[通知] 活动机会 -- 下周有个AI医疗论坛，陈大明也会去，你们可以一起去", 0),
    ("[通知] 人脉动态 -- 李娜刚完善了资料，她是CDMO领域专家，可能对你的项目有帮助", 0),
    ('用户每天来"领任务"，而不是"去搜索"', 0),
    ("人脉页面不是列表，是关系网络图", 0),
    ("私信不是自己写，AI辅助写作", 0),
])

# Slide 15: Growth
content_slide(prs, "增长策略：三个阶段", [
    ("第一阶段(0->1): 圈子里长出来", 0),
    ("  目标: 2-3个已有信任关系的校友群，KPI是用户在产品内的互动深度", 1),
    ("第二阶段(1->10): 单个模块打穿", 0),
    ("  选择供需广场 -- 价值最直接、最容易验证匹配效果", 1),
    ("  KPI: 供需匹配后开始对话的转化率", 1),
    ("第三阶段(10->100): 飞轮启动", 0),
    ("  供需产生关系 -> 关系沉淀到校友网络 -> 活动持续激活关系", 1),
    ("  -> 更强的供需匹配 -> 数据质量更高 -> 匹配更准 -> 正循环", 1),
])

# Slide 16: Summary
content_slide(prs, "核心判断与设计原则", [
    ("校友是底子，AI是引擎，活动是催化剂，供需是驱动器", 0),
    ("四个模块不是并列关系，而是一个社交飞轮", 0),
    ('原则1: 不要让用户"去找"，要让AI"来推"', 0),
    ('原则2: 每条AI推送必须附"为什么"', 0),
    ('原则3: 让下一步动作极简 -- AI推荐 -> 用户只需要"确认"或"修改"', 0),
    ("校友名录是底子，AI让底子活起来", 0),
])

# Slide 17: Next steps
content_slide(prs, "下一步建议", [
    ('先做Demo验证 -- 展示AI如何把两个有共同标签、互补供需的校友"牵线"', 0),
    ("聚焦一个模块突破 -- 建议从供需广场切入，因为价值最直接", 0),
    ("用数据说话 -- 前期不在乎用户数量，在乎匹配成功率和用户满意度", 0),
    ("给客户的核心一句话:", 0),
    ("  校友名录是底子，AI让底子活起来", 1),
    ('  不是帮你找到人，而是帮你发现你不知道该找谁的连接', 1),
])

out = "/Users/wangsheng/luojing/AI商务社交解决方案.pptx"
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
